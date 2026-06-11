import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    # Color Palette
    bg_color = RGBColor(16, 26, 33)      # Slate Dark
    text_white = RGBColor(255, 255, 255) # Title Text
    text_muted = RGBColor(160, 174, 192) # Description Text
    accent_cyan = RGBColor(79, 209, 197) # Cyan Accent
    accent_mint = RGBColor(110, 231, 183) # Mint Accent
    accent_red = RGBColor(248, 113, 113)  # Red Alert

    # Helper function to set solid background color
    def set_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    # Helper function to create clean title on slides
    def add_slide_header(slide, title_text, category="ZEROTOUCH SRE"):
        # Category label
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        cat_tf.margin_left = cat_tf.margin_top = cat_tf.margin_bottom = cat_tf.margin_right = 0
        cat_p = cat_tf.paragraphs[0]
        cat_p.text = category.upper()
        cat_p.font.name = "Arial"
        cat_p.font.size = Pt(11)
        cat_p.font.bold = True
        cat_p.font.color.rgb = accent_cyan

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.margin_left = title_tf.margin_top = title_tf.margin_bottom = title_tf.margin_right = 0
        title_p = title_tf.paragraphs[0]
        title_p.text = title_text
        title_p.font.name = "Arial"
        title_p.font.size = Pt(28)
        title_p.font.bold = True
        title_p.font.color.rgb = text_white

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Dark Theme)
    # -------------------------------------------------------------
    slide_layout = prs.slide_layouts[6] # Blank
    slide1 = prs.slides.add_slide(slide_layout)
    set_bg(slide1)

    # Subtitle Category
    box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(0.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "GOOGLE CLOUD RAPID AGENT HACKATHON  |  DYNATRACE TRACK"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = accent_cyan

    # Big Title
    box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.7), Inches(11.33), Inches(1.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "ZeroTouch SRE"
    p.font.name = "Arial"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = text_white

    # Tagline/Subtitle
    box = slide1.shapes.add_textbox(Inches(1.0), Inches(4.3), Inches(11.33), Inches(0.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Autonomous incident triage and mitigation planning for production outages."
    p.font.name = "Arial"
    p.font.size = Pt(20)
    p.font.color.rgb = text_muted

    # Subtext / Built by
    box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.33), Inches(0.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Built by Pratik Shah  •  Powered by Google Gemini & Dynatrace OpenPipeline"
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.color.rgb = accent_mint

    # -------------------------------------------------------------
    # SLIDE 2: The Problem
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(slide_layout)
    set_bg(slide2)
    add_slide_header(slide2, "The 3 AM Scramble: Production Incident Chaos")

    # Column 1: The Outage Reality
    box = slide2.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "The Challenge"
    p.font.name = "Arial"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = accent_cyan
    p.space_after = Pt(14)

    p2 = tf.add_paragraph()
    p2.text = "When critical systems go down, SRE and DevOps teams lose invaluable time doing context-switching: searching through logs, analyzing dashboards, and opening incident channels."
    p2.font.name = "Arial"
    p2.font.size = Pt(15)
    p2.font.color.rgb = text_muted
    p2.space_after = Pt(14)

    p3 = tf.add_paragraph()
    p3.text = "• First 15 Minutes are wasted purely forming a hypothesis.\n• Alert Storms create high cognitive load and noise.\n• Manual Runbooks are static, outdated, or hard to locate."
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.color.rgb = text_white
    p3.space_after = Pt(14)

    # Column 2: Pain Points
    box2 = slide2.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5))
    tf2 = box2.text_frame
    tf2.word_wrap = True

    p_r = tf2.paragraphs[0]
    p_r.text = "The Opportunity"
    p_r.font.name = "Arial"
    p_r.font.size = Pt(20)
    p_r.font.bold = True
    p_r.font.color.rgb = accent_mint
    p_r.space_after = Pt(14)

    p_r2 = tf2.add_paragraph()
    p_r2.text = "What if we could automate the entire first phase of incident response? Go from raw webhook alert to fully-diagnosed root cause, safe mitigations, and runbooks in under 5 seconds."
    p_r2.font.name = "Arial"
    p_r2.font.size = Pt(15)
    p_r2.font.color.rgb = text_muted
    p_r2.space_after = Pt(14)

    p_r3 = tf2.add_paragraph()
    p_r3.text = "• Instant Triage: Immediate diagnosis using real-time telemetry.\n• Bidirectional Audit: observability systems are kept updated.\n• Safe Execution: Gated actions to protect the production surface."
    p_r3.font.name = "Arial"
    p_r3.font.size = Pt(14)
    p_r3.font.color.rgb = text_white

    # -------------------------------------------------------------
    # SLIDE 3: The Solution / Product Showcase
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(slide_layout)
    set_bg(slide3)
    add_slide_header(slide3, "ZeroTouch SRE: Webhook-to-Artifact Automation")

    # Column 1: Core Value Proposition
    box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.2), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Autonomous Response Loop"
    p.font.name = "Arial"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = accent_cyan
    p.space_after = Pt(10)

    features = [
        ("Webhook Integration", "Receives raw alerts directly from PagerDuty, Alertmanager, or Dynatrace webhooks."),
        ("Bidirectional Telemetry", "Pushes structured audit events into Dynatrace OpenPipeline AND pulls Grail logs dynamically."),
        ("Gemini-Powered Reasoning", "Determines root cause with confidence score and suggests allowlisted mitigations."),
        ("Auto-Generated Artifacts", "Instantly writes Post-Mortem reports, JSON runbooks, and detailed execution traces.")
    ]

    for title, desc in features:
        p_title = tf.add_paragraph()
        p_title.text = f"✔ {title}"
        p_title.font.name = "Arial"
        p_title.font.size = Pt(14)
        p_title.font.bold = True
        p_title.font.color.rgb = accent_mint
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Arial"
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = text_muted
        p_desc.space_after = Pt(8)

    # Column 2: Screenshot (Landing Page)
    img_path = r"C:\Users\prati\Downloads\Projects\ZeroTouch SRE\assets\screenshots\cloud-01-landing.png"
    if os.path.exists(img_path):
        slide3.shapes.add_picture(img_path, Inches(6.4), Inches(1.8), width=Inches(6.1), height=Inches(4.6))

    # -------------------------------------------------------------
    # SLIDE 4: Live Demo / Interactive Workbench
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(slide_layout)
    set_bg(slide4)
    add_slide_header(slide4, "Live Outage Analysis: Checkout Incident")

    # Column 1: Outage Details
    box = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.2), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Incident: CPU Spike & 500 Surge"
    p.font.name = "Arial"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = accent_cyan
    p.space_after = Pt(12)

    details = [
        ("1. Root Cause Analysis", "Detected checkout-api saturated CPU at 88.0% correlated with a recent deployment. HTTP 500 error rate surged to 5.5%."),
        ("2. Safe Mitigations Planned", "Proposed: scale_service (scale checkout-api), rollback_release (revert last deployment), open_incident_channel."),
        ("3. Telemetry Evidence", "Verified using live Grail logs via Dynatrace APIs. Log context analyzed by Gemini to generate diagnosis."),
        ("4. Performance Metrics", "Triage loop completed in 4.1 seconds. Cost: < 0.05 INR using optimized Gemini token controls.")
    ]

    for title, desc in details:
        p_t = tf.add_paragraph()
        p_t.text = title
        p_t.font.name = "Arial"
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = text_white
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Arial"
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = text_muted
        p_d.space_after = Pt(8)

    # Column 2: Screenshot (Incident Result Page)
    img_path_2 = r"C:\Users\prati\Downloads\Projects\ZeroTouch SRE\assets\screenshots\cloud-02-incident-result.png"
    if os.path.exists(img_path_2):
        slide4.shapes.add_picture(img_path_2, Inches(6.4), Inches(1.8), width=Inches(6.1), height=Inches(4.6))

    # -------------------------------------------------------------
    # SLIDE 5: Safety Model (Policy Gated)
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(slide_layout)
    set_bg(slide5)
    add_slide_header(slide5, "Safety Model: Useful, Not Dangerous")

    # Column 1: Core Safety Principles
    box = slide5.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.8), Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Designed for Production Safety"
    p.font.name = "Arial"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = accent_cyan
    p.space_after = Pt(14)

    safety_points = [
        ("Simulation-Only execution by default. No live production writes.", "Allows full agent evaluation and scenario reviews without risk of service disruption."),
        ("Strict operation allowlists (scale, rollback, notify).", "Any command or script outside the allowed SRE playbook is blocked instantly at the gate."),
        ("Cost limits and token budgets checked on every single execution.", "Prevents runaway loops. Estimated token burn tracked in INR against custom budgets.")
    ]

    for title, desc in safety_points:
        p_t = tf.add_paragraph()
        p_t.text = f"🛡 {title}"
        p_t.font.name = "Arial"
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = text_white
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Arial"
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = text_muted
        p_d.space_after = Pt(12)

    # Column 2: Visual (Connect / Policy Settings Mockup)
    img_path_3 = r"C:\Users\prati\Downloads\Projects\ZeroTouch SRE\assets\screenshots\cloud-03-scenario.png"
    if os.path.exists(img_path_3):
        slide5.shapes.add_picture(img_path_3, Inches(7.0), Inches(1.8), width=Inches(5.5), height=Inches(4.6))

    # -------------------------------------------------------------
    # SLIDE 6: Architecture / Technical Stack
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(slide_layout)
    set_bg(slide6)
    add_slide_header(slide6, "Under the Hood: Cloud Architecture & Integration")

    # Column 1: Tech Stack Bullet Points
    box = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Architecture Components"
    p.font.name = "Arial"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = accent_cyan
    p.space_after = Pt(14)

    tech_points = [
        ("FastAPI Backend", "Lightweight, async Python framework serving high-speed API endpoints and incident review templates."),
        ("Google Cloud Run", "Serverless deployment scaling to zero, ensuring minimal idle cost while remaining instantly responsive to incoming webhooks."),
        ("Dynatrace Grail & OpenPipeline", "Bidirectional flow via standard REST APIs. Pushes alert audit trails and pulls real-time log evidence."),
        ("Vertex AI / Google Gemini", "High-context LLM reasoning about telemetry, generating structured incident reports and executable action plans.")
    ]

    for title, desc in tech_points:
        p_t = tf.add_paragraph()
        p_t.text = f"⚡ {title}"
        p_t.font.name = "Arial"
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = accent_mint
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Arial"
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = text_muted
        p_d.space_after = Pt(10)

    # Column 2: Architecture Loop diagram
    img_path_4 = r"C:\Users\prati\Downloads\Projects\ZeroTouch SRE\assets\screenshots\03-architecture-loop.png"
    if os.path.exists(img_path_4):
        slide6.shapes.add_picture(img_path_4, Inches(6.8), Inches(2.2), width=Inches(5.7), height=Inches(3.8))

    # Save
    out_path = r"C:\Users\prati\Downloads\Projects\ZeroTouch SRE\ZeroTouch_SRE_Presentation.pptx"
    prs.save(out_path)
    print(f"Presentation saved to {out_path}")

if __name__ == "__main__":
    create_presentation()
